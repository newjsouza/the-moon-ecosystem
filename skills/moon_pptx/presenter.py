"""PPTX creation and conversion using LibreOffice."""
import asyncio
from pathlib import Path
from core.agent_base import TaskResult


class MoonPPTXPresenter:
    """Create and convert PowerPoint presentations."""

    LIBREOFFICE_CMD = "libreoffice"

    async def create_from_slides(self, slides: list[dict], output_path: str,
                                  title: str = "The Moon Presentation", **kwargs) -> TaskResult:
        """
        Cria PPTX a partir de lista de slides.
        Cada slide: {"title": str, "content": str, "notes": str (opcional)}
        """
        start = asyncio.get_event_loop().time()
        try:
            import importlib.util
            if importlib.util.find_spec("pptx") is None:
                return TaskResult(success=False,
                                  error="python-pptx não instalado: pip install python-pptx")

            from pptx import Presentation
            from pptx.util import Inches, Pt

            prs = Presentation()
            title_slide_layout = prs.slide_layouts[0]
            content_layout = prs.slide_layouts[1]

            # Slide de título
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = title
            if len(slides) > 0 and "subtitle" in slides[0]:
                slide.placeholders[1].text = slides[0].get("subtitle", "")

            for slide_data in slides:
                s = prs.slides.add_slide(content_layout)
                s.shapes.title.text = slide_data.get("title", "")
                if s.placeholders[1]:
                    s.placeholders[1].text = slide_data.get("content", "")
                if slide_data.get("notes"):
                    s.notes_slide.notes_text_frame.text = slide_data["notes"]

            prs.save(output_path)
            path = Path(output_path)

            return TaskResult(
                success=True,
                data={"pptx_path": output_path, "slides": len(slides),
                      "size_bytes": path.stat().st_size},
                execution_time=asyncio.get_event_loop().time() - start
            )
        except Exception as e:
            return TaskResult(success=False, error=str(e))

    async def convert_to_pdf(self, source_path: str,
                              output_dir: str = None, **kwargs) -> TaskResult:
        """Exporta PPTX para PDF via LibreOffice (para relatórios OmniChannelStrategist)."""
        start = asyncio.get_event_loop().time()
        try:
            src = Path(source_path)
            if not src.exists():
                return TaskResult(success=False, error=f"Arquivo não encontrado: {source_path}")

            out_dir = output_dir or str(src.parent)
            cmd = [self.LIBREOFFICE_CMD, "--headless", "--convert-to", "pdf",
                   "--outdir", out_dir, str(src)]

            proc = await asyncio.create_subprocess_exec(
                *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
            )
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=60)

            if proc.returncode != 0:
                return TaskResult(success=False,
                                  error=f"LibreOffice erro: {stderr.decode().strip()}")

            pdf_path = Path(out_dir) / (src.stem + ".pdf")
            return TaskResult(
                success=True,
                data={"pdf_path": str(pdf_path)},
                execution_time=asyncio.get_event_loop().time() - start
            )
        except Exception as e:
            return TaskResult(success=False, error=str(e))